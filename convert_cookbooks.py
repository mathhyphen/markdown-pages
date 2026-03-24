#!/usr/bin/env python3
"""Convert Claude Cookbooks markdown files to HTML and PPTX slides."""

import os
import re
import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC_DIR = "D:/apps/markdown-pages/cookbooks"
HTML_DIR = "D:/apps/markdown-pages/cookbooks"
PPTX_DIR = "D:/apps/markdown-pages/cookbooks-slides"

def markdown_to_html(md_content):
    """Convert markdown to HTML with proper styling."""
    html_template = '''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; max-width: 900px; margin: 0 auto; padding: 2rem; line-height: 1.6; }}
        h1 {{ color: #1a1a2e; border-bottom: 3px solid #4a90d9; padding-bottom: 0.5rem; }}
        h2 {{ color: #16213e; margin-top: 2rem; }}
        h3 {{ color: #0f3460; }}
        code {{ background: #f4f4f4; padding: 0.2rem 0.4rem; border-radius: 3px; }}
        pre {{ background: #1e1e1e; color: #d4d4d4; padding: 1rem; border-radius: 8px; overflow-x: auto; }}
        blockquote {{ border-left: 4px solid #4a90d9; margin: 1rem 0; padding-left: 1rem; color: #666; }}
        a {{ color: #4a90d9; }}
        table {{ border-collapse: collapse; width: 100%; }}
        th, td {{ border: 1px solid #ddd; padding: 0.5rem; text-align: left; }}
        th {{ background: #f4f4f4; }}
    </style>
</head>
<body>
{content}
</body>
</html>'''
    title = re.search(r'^#\s+(.+)$', md_content, re.MULTILINE)
    title = title.group(1) if title else "Claude Cookbooks"
    html_content = markdown.markdown(md_content, extensions=['tables', 'fenced_code', 'codehilite'])
    return html_template.format(title=title, content=html_content)

def extract_title_and_content(md_content):
    """Extract title and content sections from markdown."""
    lines = md_content.split('\n')
    title = ""
    content_blocks = []

    for line in lines:
        if line.startswith('# ') and not title:
            title = line[2:].strip()
        elif line.startswith('## '):
            content_blocks.append(('section', line[3:].strip()))
        elif line.startswith('### '):
            content_blocks.append(('subsection', line[4:].strip()))
        elif line.startswith('- ') or line.startswith('* '):
            content_blocks.append(('bullet', line[2:].strip()))
        elif line.strip() and content_blocks:
            last = content_blocks[-1]
            if last[0] == 'text':
                content_blocks[-1] = ('text', last[1] + ' ' + line.strip())
            else:
                content_blocks.append(('text', line.strip()))

    return title, content_blocks

def create_pptx(md_file, md_content):
    """Create PPTX presentation from markdown content."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title, content_blocks = extract_title_and_content(md_content)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11.333)
    height = Inches(2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 26, 46)
    p.alignment = PP_ALIGN.CENTER

    # Content slides
    current_text = []
    current_header = "内容"

    def add_content_slide():
        if not current_text:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Header
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = current_header
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(74, 144, 217)

        # Content
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, text in enumerate(current_text[:15]):  # Max 15 bullets
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + text[:100] + ("..." if len(text) > 100 else "")
            p.font.size = Pt(18)
            p.space_after = Pt(8)

    for block_type, block_content in content_blocks:
        if block_type in ('section', 'subsection'):
            if current_text:
                add_content_slide()
                current_text = []
            current_header = block_content
        elif block_type == 'bullet':
            current_text.append(block_content)
        elif block_type == 'text' and current_text:
            current_text[-1] += " " + block_content

    if current_text:
        add_content_slide()

    # Save
    base = os.path.splitext(os.path.basename(md_file))[0]
    output_path = os.path.join(PPTX_DIR, f"{base}.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")

def main():
    os.makedirs(HTML_DIR, exist_ok=True)
    os.makedirs(PPTX_DIR, exist_ok=True)

    for filename in sorted(os.listdir(SRC_DIR)):
        if not filename.endswith('.md'):
            continue

        filepath = os.path.join(SRC_DIR, filename)
        with open(filepath, 'r', encoding='utf-8') as f:
            md_content = f.read()

        # Create HTML
        html_content = markdown_to_html(md_content)
        base = os.path.splitext(filename)[0]
        html_path = os.path.join(HTML_DIR, f"{base}.html")
        with open(html_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        print(f"Created: {html_path}")

        # Create PPTX
        create_pptx(filepath, md_content)

    print("\nDone!")

if __name__ == "__main__":
    main()
